from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

W, H = 40, 30
prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])

CORNELL_RED = RGBColor(0xB3, 0x1B, 0x1B)
DARK_RED = RGBColor(0x7B, 0x11, 0x13)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BLUE = RGBColor(0x15, 0x65, 0xC0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE1)
BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
BOX_BG = RGBColor(0xFF, 0xFF, 0xFF)
RED_TEXT = RGBColor(0xC6, 0x28, 0x28)

M = 0.3
COL_GAP = 0.2
ROW_GAP = 0.15
HEADER_H = 2.0
FOOTER_H = 0.6

n_cols = 4
col_w = (W - 2*M - (n_cols-1)*COL_GAP) / n_cols
content_top = M + HEADER_H + ROW_GAP
content_h = H - M - FOOTER_H - content_top - ROW_GAP

def col_x(i):
    return M + i * (col_w + COL_GAP)

def add_rect(slide, left, top, width, height, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or BOX_BG
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False

    shape.adjustments[0] = 0.02
    return shape

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def set_text(tf, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return p

def add_para(tf, text, size=14, bold=False, color=BLACK, space_before=0, space_after=0, bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet:
        p.level = 0
    return p

def add_image(slide, path, left, top, width):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))

header = add_rect(slide, M, M, W - 2*M, HEADER_H, fill=CORNELL_RED)

tb = add_textbox(slide, M + 0.3, M + 0.2, W - 2*M - 8, 1.2)
set_text(tb.text_frame, 'Inference Optimization for Hybrid Diffusion-Autoregressive\nLanguage Models: Profiling, Fusion, and Framework Overhead',
         size=32, bold=True, color=WHITE)

tb2 = add_textbox(slide, M + 0.3, M + 1.3, W - 2*M - 8, 0.5)
set_text(tb2.text_frame, 'Srivatsa Kundurthy    |    Cornell University    |    CS 5220: Applied High-Performance and Parallel Computing    |    Spring 2026',
         size=16, color=WHITE)

tb3 = add_textbox(slide, W - M - 7.5, M + 0.3, 7, 1.4)
tf3 = tb3.text_frame
tf3.word_wrap = True
set_text(tf3, 'STAR-LDM (COLM 2025)', size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
add_para(tf3, 'Lovelace, Belardi, Zalouk, Polavaram,', size=13, color=WHITE)
add_para(tf3, 'Kundurthy, Weinberger', size=13, bold=True, color=WHITE)
add_para(tf3, 'Apple M4 Max  |  128 GB  |  Metal GPU', size=13, color=WHITE, space_before=6)

x1 = col_x(0)
y = content_top

b1 = add_rect(slide, x1, y, col_w, 7.5, fill=BOX_BG, border=CORNELL_RED)
tb = add_textbox(slide, x1 + 0.15, y + 0.1, col_w - 0.3, 7.2)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, '1. STAR-LDM: Hybrid Diffusion + Autoregressive', size=20, bold=True, color=CORNELL_RED)
add_para(tf, '', size=6)
add_para(tf, 'STAR-LDM combines 50-step latent diffusion planning with GPT-2 Large autoregressive decoding, achieving dramatically better text quality:', size=14, space_after=4)
add_para(tf, '', size=4)
add_para(tf, 'Model              Params    MAUVE ↑', size=13, bold=True, color=BLACK, space_after=2)
add_para(tf, 'GPT-2 Large      770M       85.2', size=13, space_after=1)
add_para(tf, 'GPT-2 XL           1.5B        86.6', size=13, space_after=1)
add_para(tf, 'STAR-LDM         956M       94.6  (+8 points)', size=14, bold=True, color=GREEN, space_after=4)
add_para(tf, '', size=6)
add_para(tf, 'The Stop → Think → AutoRegress Pipeline', size=17, bold=True, color=BLACK, space_before=4)
add_para(tf, '', size=4)
add_para(tf, '1. STOP: Encode prefix through GPT-2, cache KV projections.', size=13, space_after=2)
add_para(tf, '2. THINK: 50-step diffusion loop. Each step:', size=13, space_after=2)
add_para(tf, '   z_t (768d) → SPG (6 layers) → 8 soft prompts (1280d)', size=12, color=BLUE, space_after=1)
add_para(tf, '   → GPT-2 Large (36 layers, 770M) → hidden states', size=12, color=RED_TEXT, space_after=1)
add_para(tf, '   → ScoreNet (6 layers) → v-prediction → DDPM update → z_{t-1}', size=12, color=GREEN, space_after=4)
add_para(tf, '3. AUTOREGRESS: Final soft prompts condition GPT-2 for token generation.', size=13, space_after=4)
add_para(tf, '', size=6)
add_para(tf, 'The Problem', size=17, bold=True, color=RED_TEXT, space_before=4)
add_para(tf, 'Each step invokes the full 770M-param GPT-2 backbone. At 50 steps, STAR-LDM is 2.3x slower than GPT-2 Large. The authors note their implementation is "not optimized for inference speed."', size=13, space_after=2)

y2 = y + 7.7
add_rect(slide, x1, y2, col_w, col_w * 0.6 + 0.3, fill=BOX_BG)
tb = add_textbox(slide, x1 + 0.15, y2 + 0.1, col_w - 0.3, 0.3)
set_text(tb.text_frame, 'Architecture: One Diffusion Step', size=16, bold=True, color=CORNELL_RED)
add_image(slide, 'figures/fig0_architecture.png', x1 + 0.15, y2 + 0.4, col_w - 0.3)

x2 = col_x(1)

b2 = add_rect(slide, x2, y, col_w, 5.5, fill=BOX_BG, border=CORNELL_RED)
tb = add_textbox(slide, x2 + 0.15, y + 0.1, col_w - 0.3, 0.3)
set_text(tb.text_frame, '2. Profiling: Where Does Time Go?', size=20, bold=True, color=CORNELL_RED)
add_image(slide, 'figures/fig1_profiling_pie.png', x2 + 0.2, y + 0.5, col_w - 0.4)
tb = add_textbox(slide, x2 + 0.15, y + 4.2, col_w - 0.3, 1.2)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, '50 C4 prompts, M4 Max, 50 steps, unoptimized.', size=12, color=BLACK)
add_para(tf, 'GPT-2 = 77.2% of total inference time.', size=14, bold=True, color=RED_TEXT)

y2b = y + 5.7
b2b = add_rect(slide, x2, y2b, col_w, 4.8, fill=BOX_BG, border=CORNELL_RED)
tb = add_textbox(slide, x2 + 0.15, y2b + 0.1, col_w - 0.3, 0.3)
set_text(tb.text_frame, 'The Real Bottleneck: Framework Overhead', size=17, bold=True, color=CORNELL_RED)
add_image(slide, 'figures/fig4_dispatch_count.png', x2 + 0.3, y2b + 0.5, col_w - 0.6)
tb = add_textbox(slide, x2 + 0.15, y2b + 3.6, col_w - 0.3, 1.0)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, '6,185 ATen dispatches per call; only 432 (7%) compute.', size=13, bold=True, color=BLACK)
add_para(tf, 'The rest: mask creation, cache management, tensor reshaping. At ~3.4μs each, dispatch overhead alone = ~21ms/call.', size=12)

y2c = y2b + 5.0
b2c = add_rect(slide, x2, y2c, col_w, col_w * 0.75 + 0.8, fill=BOX_BG, border=CORNELL_RED)
tb = add_textbox(slide, x2 + 0.15, y2c + 0.1, col_w - 0.3, 0.3)
set_text(tb.text_frame, 'Roofline: Dispatch-Latency-Bound', size=17, bold=True, color=CORNELL_RED)
add_image(slide, 'figures/fig3_roofline.png', x2 + 0.15, y2c + 0.45, col_w - 0.3)
tb = add_textbox(slide, x2 + 0.15, y2c + col_w*0.75 + 0.1, col_w - 0.3, 0.5)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, 'Micro-transformer ops sit at <1% of the roofline ceiling: dispatch-latency-bound rather than compute- or memory-bound.', size=12, bold=True, color=RED_TEXT)

x3 = col_x(2)

b3 = add_rect(slide, x3, y, col_w, 8.0, fill=BOX_BG, border=CORNELL_RED)
tb = add_textbox(slide, x3 + 0.15, y + 0.1, col_w - 0.3, 7.8)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, '3. Optimizations', size=20, bold=True, color=CORNELL_RED)
add_para(tf, '', size=6)
add_para(tf, 'A. Streamlined GPT-2 Forward', size=16, bold=True, color=BLACK, space_before=4)
add_para(tf, 'Replace HuggingFace GPT2LMHeadModel.forward() with 40 lines of PyTorch performing only essential ops. Eliminates mask construction, cache management, tensor format conversions.', size=13, space_after=4)
add_para(tf, '6,185 → 432 dispatches  |  3.6x faster per-call', size=15, bold=True, color=GREEN, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'B. KV-Cache with Pre-Allocated Buffers', size=16, bold=True, color=BLACK, space_before=4)
add_para(tf, 'Compute prefix KV once, reuse across all 50 steps. Pre-allocating contiguous buffers eliminates 1,800 tensor allocations per generation.', size=13, space_after=2)
add_para(tf, 'Per-step cost: O(1) vs O(prefix_len)', size=14, bold=True, color=GREEN, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'C. Metal Compute Shaders (1,700 lines)', size=16, bold=True, color=BLACK, space_before=4)
add_para(tf, '6 custom Apple Metal GPU kernels:', size=13, space_after=2)
add_para(tf, '✓ RMSNorm+FiLM:  3→1 dispatch, 1.69x', size=13, color=GREEN, space_after=1)
add_para(tf, '✓ Tiny Attention:  register-resident 8×8, 2.39x', size=13, color=GREEN, space_after=1)
add_para(tf, '✗ DDPM Step:  dispatch > compute, 0.35x', size=13, color=RED_TEXT, space_after=1)
add_para(tf, '✗ Fused FFN:  can\'t beat Apple BLAS, 0.01x', size=13, color=RED_TEXT, space_after=1)
add_para(tf, '✗ Decode-N Attn:  wins only at KV<40', size=13, color=RED_TEXT, space_after=1)
add_para(tf, '✗ Spec Verify:  dispatch overhead, 0.13x', size=13, color=RED_TEXT, space_after=2)
add_para(tf, '2 of 6 kernels deliver speedups.', size=13, bold=True, space_after=2)

y3b = y + 8.2
b3b = add_rect(slide, x3, y3b, col_w, 7.5, fill=BOX_BG, border=CORNELL_RED)
tb = add_textbox(slide, x3 + 0.15, y3b + 0.1, col_w - 0.3, 0.3)
set_text(tb.text_frame, '4. Results: Prefix Length Scaling', size=20, bold=True, color=CORNELL_RED)
add_image(slide, 'figures/fig2_prefix_scaling.png', x3 + 0.15, y3b + 0.5, col_w - 0.3)
tb = add_textbox(slide, x3 + 0.15, y3b + 3.8, col_w - 0.3, 3.5)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, 'Prefix    Baseline    Optimized    Speedup', size=13, bold=True)
add_para(tf, '  16        921 ms      748 ms        1.23x', size=12, space_after=1)
add_para(tf, '  64      1,072 ms      765 ms        1.40x', size=12, space_after=1)
add_para(tf, '128      1,321 ms      801 ms        1.65x', size=12, space_after=1)
add_para(tf, '256      1,835 ms      853 ms        2.15x', size=12, space_after=1)
add_para(tf, '512      2,894 ms      982 ms        2.95x', size=13, bold=True, color=GREEN, space_after=4)
add_para(tf, '', size=4)
add_para(tf, '100-Prompt Evaluation (20 steps, prefixes 5–200 tok):', size=14, bold=True, space_before=4)
add_para(tf, 'Baseline: 1,590 ms (std 279)  →  Optimized: 1,379 ms (std 50)', size=13, space_after=1)
add_para(tf, '1.15x speedup, 5.6x lower variance', size=14, bold=True, color=GREEN)

x4 = col_x(3)

b4 = add_rect(slide, x4, y, col_w, 9.5, fill=BLUE_BG, border=BLUE)
tb = add_textbox(slide, x4 + 0.15, y + 0.1, col_w - 0.3, 9.2)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, '5. Future Work (for Final Report)', size=20, bold=True, color=BLUE)
add_para(tf, '', size=6)
add_para(tf, 'A. Cross-Platform CUDA Comparison', size=16, bold=True, color=BLACK, space_before=6)
add_para(tf, 'Run identical benchmarks on NVIDIA A100. Test whether framework overhead is MPS-specific or general. Does torch.compile with CUDA Graphs succeed where MPS failed?', size=13, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'B. Quality Validation', size=16, bold=True, color=BLACK, space_before=4)
add_para(tf, 'Generate 5,000 C4 continuations with both pipelines. Compute MAUVE and perplexity to verify optimizations preserve output quality. Critical: our fast forward drops causal masking on the 8 soft prompts.', size=13, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'C. Model Scheduling', size=16, bold=True, color=BLACK, space_before=4)
add_para(tf, 'Profile v-prediction error at each noise level. At early/late steps, GPT-2 may add minimal value, allowing replacement with a cheap SPG+ScoreNet-only path (~3ms vs ~12ms/step). Training-free, potential 30 to 40% further speedup.', size=13, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'D. Consistency Distillation (Stretch)', size=16, bold=True, color=BLACK, space_before=4)
add_para(tf, 'Train student score network to map z_t → z_0 in 1 to 4 steps (LCM-style). STAR-LDM\'s 768-dim space is smoother than pixel space, so few-step distillation should be effective.', size=13, space_after=2)
add_para(tf, 'At 4 steps: ~550ms total (faster than GPT-2 Large)', size=14, bold=True, color=GREEN, space_after=2)

y4b = y + 9.7
b4b = add_rect(slide, x4, y4b, col_w, 5.5, fill=AMBER_BG, border=ORANGE)
tb = add_textbox(slide, x4 + 0.15, y4b + 0.1, col_w - 0.3, 5.3)
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, 'Key Takeaway', size=20, bold=True, color=ORANGE)
add_para(tf, '', size=6)
add_para(tf, 'The primary bottleneck in hybrid diffusion-AR inference is framework dispatch overhead, not GPU kernel performance.', size=15, bold=True, color=BLACK, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'HuggingFace issues 14x more operator dispatches than needed. A 40-line PyTorch rewrite gives a larger speedup (3.6x on 70% of pipeline) than six custom Metal GPU kernels combined.', size=13, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'Roofline analysis reveals a dispatch-latency-bound regime where micro-transformer ops achieve <1% of the hardware ceiling: a third regime beyond compute-bound and memory-bound.', size=13, space_after=6)
add_para(tf, '', size=4)
add_para(tf, 'Negative results: Picard parallel diffusion (convergence too slow), torch.compile on MPS (2–3x slower), fused FFN kernel (127x slower than BLAS).', size=12, space_after=4)
add_para(tf, '', size=6)
add_para(tf, 'References', size=14, bold=True, color=ORANGE, space_before=4)
add_para(tf, 'Lovelace et al., COLM 2025 • Dao et al., NeurIPS 2022 • Aminabadi et al., SC 2022 • Williams et al., CACM 2009 • Yuan et al., 2024 • Ansel et al., ASPLOS 2024', size=10, space_after=2)

footer = add_rect(slide, M, H - M - FOOTER_H, W - 2*M, FOOTER_H, fill=RGBColor(0x33, 0x33, 0x33))
tb = add_textbox(slide, M + 0.3, H - M - FOOTER_H + 0.1, W/2, 0.4)
set_text(tb.text_frame, 'Code: 1,700 lines Metal/Obj-C++ (6 GPU kernels) + optimized PyTorch pipeline  |  Eval: 100 prompts, 4 prefix-length buckets',
         size=12, color=WHITE)
tb = add_textbox(slide, W/2, H - M - FOOTER_H + 0.1, W/2 - M - 0.3, 0.4)
set_text(tb.text_frame, 'CS 5220: Applied High-Performance and Parallel Computing  |  Prof. Giulia Guidi  |  Cornell University  |  Spring 2026',
         size=12, color=WHITE, align=PP_ALIGN.RIGHT)

out_path = 'poster.pptx'
prs.save(out_path)
print(f'Saved to {out_path}')
print(f'Slide size: {W}x{H} inches')
